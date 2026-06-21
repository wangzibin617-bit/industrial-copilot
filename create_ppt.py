"""
Industrial Copilot — AI PM Interview Presentation
7 slides with screenshots, optimized for AI Product Manager narrative
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── CONFIG ──
SCREENSHOT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "screenshots")
OUTPUT_PATH   = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Industrial_Copilot_Interview.pptx")

# ── PALETTE ──
NAVY       = RGBColor(0x0B, 0x15, 0x2E)
DARK_CARD  = RGBColor(0x13, 0x20, 0x42)
STEEL      = RGBColor(0x2D, 0x5F, 0x8A)
TEAL       = RGBColor(0x00, 0xBA, 0xB3)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
DIM_GRAY   = RGBColor(0x64, 0x72, 0x82)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
RED        = RGBColor(0xEF, 0x44, 0x44)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── HELPERS ──

def bg(slide, c=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c

def rect(slide, l, t, w, h, fill=None, border=None, corner=False):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    return s

def tb(slide, l, t, w, h, text="", size=16, color=WHITE, bold=False,
       align=PP_ALIGN.LEFT, font="Microsoft YaHei", spacing=1.3, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = Pt(size * spacing)
    return box

def rich(slide, l, t, w, h, paras):
    """paras: list of dicts {text, size, color, bold, align, spacing, space_before}"""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, pdata in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = pdata.get("text", "")
        p.font.size = Pt(pdata.get("size", 14))
        p.font.color.rgb = pdata.get("color", WHITE)
        p.font.bold = pdata.get("bold", False)
        p.font.name = pdata.get("font", "Microsoft YaHei")
        p.alignment = pdata.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(pdata.get("spacing", 4))
        p.space_before = Pt(pdata.get("space_before", 0))
        fs = pdata.get("size", 14)
        p.line_spacing = Pt(fs * pdata.get("line_spacing", 1.25))
    return box

def section_header(slide, title, subtitle=None):
    rect(slide, Inches(0.7), Inches(0.45), Inches(0.5), Pt(4), fill=TEAL)
    tb(slide, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.55), title, size=30, bold=True)
    if subtitle:
        tb(slide, Inches(0.7), Inches(1.05), Inches(11.5), Inches(0.35), subtitle, size=12, color=MID_GRAY)

def page_num(slide, n):
    tb(slide, Inches(12.3), Inches(7.08), Inches(0.8), Inches(0.3), str(n), size=9, color=DIM_GRAY, align=PP_ALIGN.RIGHT)

def accent_bar(slide, l, t, w):
    rect(slide, l, t, w, Pt(3), fill=TEAL)

def add_screenshot(slide, filename, l, t, w, h=None):
    """Embed a screenshot image."""
    filepath = os.path.join(SCREENSHOT_DIR, filename)
    if os.path.exists(filepath):
        if h is None:
            h = w * 0.625  # 16:10 default
        return slide.shapes.add_picture(filepath, l, t, w, h)
    else:
        # Placeholder if image missing
        r = rect(slide, l, t, w, h, fill=DARK_CARD, border=STEEL)
        tb(slide, l + Inches(0.2), t + h / 2 - Inches(0.2), w - Inches(0.4), Inches(0.4),
           f"[Screenshot: {filename}]", size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
        return r

def pill_card(slide, l, t, w, h, icon, title, desc, accent=TEAL):
    """Small metric / highlight card."""
    card = rect(slide, l, t, w, h, fill=DARK_CARD, corner=True)
    rect(slide, l, t, w, Pt(3), fill=accent)
    tb(slide, l + Inches(0.15), t + Inches(0.1), w - Inches(0.3), Inches(0.3),
       f"{icon}  {title}", size=12, color=accent, bold=True)
    tb(slide, l + Inches(0.15), t + Inches(0.45), w - Inches(0.3), h - Inches(0.55),
       desc, size=10, color=MID_GRAY, spacing=1.4)

# ═══════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

# Top + bottom accent lines
rect(s, Inches(0), Inches(0), W, Pt(6), fill=TEAL)
rect(s, Inches(0), Inches(7.44), W, Pt(6), fill=TEAL)

# Left accent bar behind title
rect(s, Inches(0.9), Inches(2.0), Pt(5), Inches(3.2), fill=TEAL)

tb(s, Inches(1.25), Inches(2.15), Inches(10), Inches(0.9),
   "Industrial Copilot", size=56, bold=True)

tb(s, Inches(1.25), Inches(3.0), Inches(10), Inches(0.5),
   "AI 驱动的工业自动化销售助手", size=22, color=TEAL)

tb(s, Inches(1.25), Inches(3.55), Inches(10), Inches(0.45),
   "从发现痛点到交付 AI SaaS 产品 — 一个 AI 产品经理的 0→1 实践", size=13, color=MID_GRAY)

# Info strip
rect(s, Inches(0.9), Inches(5.6), Inches(11.5), Pt(0.5), fill=STEEL)
rich(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(1.0), [
    {"text": "AI 产品经理面试作品  ·  2026", "size": 14, "color": MID_GRAY},
    {"text": "全栈独立开发  |  Next.js 15 + Supabase + DeepSeek + Vercel AI SDK  |  2 周 MVP 交付", "size": 10, "color": DIM_GRAY, "spacing": 4},
])

page_num(s, 1)

# ═══════════════════════════════════════════
# SLIDE 2 — PAIN POINTS & PRODUCT INSIGHT
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "用户痛点 & 产品洞察", "发现真问题，定义对的产品")

# 4 pain point cards
pains = [
    ("⏳", "选型等待 1-2 天", "销售 100% 依赖售前工程师\n销售只能等，客户在流失", RGBColor(0xE8, 0x5D, 0x5D)),
    ("🔄", "需求传递失真", "口头转述 → 反复追问 ≥ 3 轮\n平均沟通损耗 30-60 min/单", AMBER),
    ("📉", "经验无法沉淀", "老销售离职 = 知识清零\n新人 3 个月才能独立接单", RGBColor(0x63, 0x9F, 0xEB)),
    ("📚", "跨品牌对比困难", "施耐德、信捷手册格式各异\n手工 Excel 对比 30-60 min/次", TEAL),
]
for i, (icon, title, desc, accent) in enumerate(pains):
    x = Inches(0.7 + i * 3.08)
    y = Inches(1.8)
    pill_card(s, x, y, Inches(2.85), Inches(1.9), icon, title, desc, accent)

# "Why AI?" insight box
rect(s, Inches(0.7), Inches(4.15), Inches(11.9), Pt(0.5), fill=STEEL)
rich(s, Inches(0.7), Inches(4.35), Inches(5.8), Inches(2.6), [
    {"text": "💡  为什么这是 AI 产品机会？", "size": 17, "color": TEAL, "bold": True, "spacing": 10},
    {"text": "工业自动化选型是规则密集 + 知识密集领域：", "size": 12, "color": LIGHT_GRAY, "spacing": 6},
    {"text": "▸ 产品参数有结构化边界（I/O、通讯协议、功率）→ 可被检索", "size": 11, "color": WHITE, "spacing": 3},
    {"text": "▸ 销售知识可沉淀（选型指南、对比话术、场景方案）→ 可被注入 Prompt", "size": 11, "color": WHITE, "spacing": 3},
    {"text": "▸ LLM 擅长理解自然语言需求 + 匹配结构化数据 → 天然方案", "size": 11, "color": WHITE, "spacing": 3},
    {"text": "", "size": 4, "color": WHITE},
    {"text": "🎯 产品目标：让不懂技术的销售，用自然语言 5 分钟内获得基于真实产品数据的推荐型号", "size": 12, "color": TEAL, "bold": True, "spacing": 2},
])

# Right side: User persona
rect(s, Inches(6.8), Inches(4.35), Inches(5.8), Inches(2.6), fill=DARK_CARD, corner=True)
rich(s, Inches(7.0), Inches(4.5), Inches(5.4), Inches(2.3), [
    {"text": "👤  MVP 目标用户", "size": 15, "color": TEAL, "bold": True, "spacing": 8},
    {"text": "张伟  |  32 岁  |  工业自动化销售 5 年", "size": 13, "color": WHITE, "bold": True, "spacing": 4},
    {"text": "", "size": 4, "color": WHITE},
    {"text": "日常：客户微信发需求 → 翻手册 → 问售前 → 等 1-2 天 → 回复客户", "size": 11, "color": MID_GRAY, "spacing": 3},
    {"text": "痛点：不懂 PLC 型号差异，选型 100% 依赖售前", "size": 11, "color": RED, "spacing": 3},
    {"text": "", "size": 4, "color": WHITE},
    {"text": "使用 AI 后：客户发需求 → 打开 Copilot 输入 → 5 分钟拿推荐 → 即刻回复", "size": 11, "color": GREEN, "spacing": 3},
    {"text": "价值：从「等售前」变成「自助选型」，响应速度提升 100x", "size": 11, "color": GREEN, "spacing": 3},
])

page_num(s, 2)

# ═══════════════════════════════════════════
# SLIDE 3 — PRODUCT DESIGN / WORKFLOW + SCREENSHOT
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "产品方案 & 核心工作流", "自然语言输入 → 数据检索 → AI 推理 → 结构化输出 → 持久化")

# Left: Flow diagram
flow_y = Inches(1.7)
steps_data = [
    ("1", "用户输入\n自然语言需求", RGBColor(0x15, 0x28, 0x50)),
    ("2", "关键词提取\nSQL ILIKE 检索", STEEL),
    ("3", "匹配产品库\n+ 销售知识库", RGBColor(0x1E, 0x6B, 0x8A)),
    ("4", "DeepSeek\n流式 AI 推理", RGBColor(0x00, 0x7A, 0x75)),
    ("5", "6 段式结构化\n回答输出", TEAL),
    ("6", "对话 & 方案\n自动持久化", RGBColor(0x00, 0x5F, 0x5A)),
]
bx, bh = Inches(1.55), Inches(1.1)
gap = Inches(0.15)

for i, (num, label, clr) in enumerate(steps_data):
    x = Inches(0.7 + i * (bx + gap))
    shape = rect(s, x, flow_y, bx, bh, fill=clr, corner=True)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{num}\n{label.split(chr(10))[0]}\n{label.split(chr(10))[1] if chr(10) in label else ''}"
    # Simpler: just put num big + label small
    tf.clear()
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label.split("\n")[0] + ("\n" + label.split("\n")[1] if len(label.split("\n")) > 1 else "")
    p2.font.size = Pt(9)
    p2.font.color.rgb = WHITE
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER
    if i < len(steps_data) - 1:
        # arrow
        ax = x + bx + Inches(0.02)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, flow_y + bh / 2 - Inches(0.08), Inches(0.11), Inches(0.16))
        arr.fill.solid(); arr.fill.fore_color.rgb = TEAL; arr.line.fill.background()

# Below: details + screenshot side by side
# Left detail column
rich(s, Inches(0.7), Inches(3.2), Inches(6.0), Inches(3.8), [
    {"text": "📋  RAG 检索 + AI 推理细节", "size": 15, "color": TEAL, "bold": True, "spacing": 8},
    {"text": "① Keyword Extraction — 中文分词 + 停用词过滤，免外部 NLP 依赖", "size": 10.5, "color": LIGHT_GRAY, "spacing": 3},
    {"text": "② SQL ILIKE — 匹配 product name / model / series / description", "size": 10.5, "color": LIGHT_GRAY, "spacing": 3},
    {"text": "③ buildDataContext — 检索结果 + source_note 注入 System Prompt", "size": 10.5, "color": LIGHT_GRAY, "spacing": 3},
    {"text": "④ streamText → Vercel AI SDK → Token-by-token 流式渲染", "size": 10.5, "color": LIGHT_GRAY, "spacing": 3},
    {"text": "⑤ onFinish 回调 — 服务端写入 messages，不依赖客户端可靠性", "size": 10.5, "color": LIGHT_GRAY, "spacing": 3},
    {"text": "", "size": 6, "color": WHITE},
    {"text": "🗣️  AI 输出 6 段式结构", "size": 15, "color": TEAL, "bold": True, "spacing": 6},
    {"text": "需求理解 → 推荐产品表格(型号+参数) → 推荐原因 → 替代方案 → 风险提示 → 参考来源", "size": 10.5, "color": WHITE, "spacing": 3},
    {"text": "每一条产品数据含 source_type + source_note，AI 回答可追溯到具体数据源", "size": 10, "color": DIM_GRAY, "spacing": 2},
])

# Right: chat screenshot
add_screenshot(s, "02_chat.png", Inches(7.0), Inches(3.2), Inches(5.7), Inches(3.8))

page_num(s, 3)

# ═══════════════════════════════════════════
# SLIDE 4 — AI PRODUCT DESIGN DECISIONS
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "AI 产品设计决策", "MVP 范围取舍  ·  System Prompt 工程  ·  防幻觉机制  ·  数据可追溯")

# Build vs Don't Build
tb(s, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.35),
   "✅  MVP 做了什么（核心假设验证）", size=15, color=GREEN, bold=True)
tb(s, Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.35),
   "❌  刻意不做什么（产品判断力）", size=15, color=RED, bold=True)

build = [
    "单用户系统 → 先验证 AI 选型价值，后扩展协作",
    "SQL ILIKE 关键词检索 → 20 条数据无需向量库",
    "EAV 属性模型 → 灵活适配不同品类参数差异",
    "服务端 onFinish 持久化对话 → 消息不丢失",
    "双认证模式 (Cookie + Bearer) → 兼顾浏览器和 API",
]
skip = [
    "多层 RBAC → 单用户无协作场景，过度设计",
    "pgvector 向量检索 → 数据量小，ILIKE 已够用",
    "Word 方案导出 → 先验证选型准确率再扩展",
    "价格/报价引擎 → 代理商进价敏感，MVP 不做定价",
    "微信小程序 → PC Web 先验证闭环",
]

for i, item in enumerate(build):
    tb(s, Inches(0.8), Inches(2.2 + i * 0.42), Inches(5.3), Inches(0.4),
       f"▸ {item}", size=10.5, color=LIGHT_GRAY)
for i, item in enumerate(skip):
    tb(s, Inches(6.9), Inches(2.2 + i * 0.42), Inches(5.3), Inches(0.4),
       f"▸ {item}", size=10.5, color=MID_GRAY)

rect(s, Inches(0.7), Inches(4.45), Inches(11.9), Pt(0.5), fill=STEEL)

# System Prompt Engineering
rich(s, Inches(0.7), Inches(4.6), Inches(5.8), Inches(2.5), [
    {"text": "📝  System Prompt 设计原则", "size": 15, "color": TEAL, "bold": True, "spacing": 8},
    {"text": "1. 角色锚定：工业自动化售前工程师，服务施耐德+信捷双品牌代理商", "size": 10.5, "color": WHITE, "spacing": 3},
    {"text": "2. 输出约束：强制 6 段式格式，确保回答质量一致", "size": 10.5, "color": WHITE, "spacing": 3},
    {"text": "3. 幻觉防护：仅基于检索数据推荐，无匹配时明确声明，不编造参数", "size": 10.5, "color": WHITE, "spacing": 3},
    {"text": "4. 角色感知：根据用户角色调整语言风格（销售用非技术语言）", "size": 10.5, "color": WHITE, "spacing": 3},
    {"text": "5. 数据溯源：每条推荐标注 source_note，告知数据可信度", "size": 10.5, "color": WHITE, "spacing": 3},
])

# Anti-hallucination 4 layers
rich(s, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.5), [
    {"text": "🛡️  四层防幻觉机制（生产级设计）", "size": 15, "color": TEAL, "bold": True, "spacing": 8},
    {"text": "", "size": 3, "color": WHITE},
])
shield_data = [
    ("① Prompt 约束", "System Prompt 明确「仅基于检索数据推荐，不编造参数」"),
    ("② 数据注入", "DB 数据作为上下文边界，AI 推理被约束在此范围内"),
    ("③ 来源标注", "每条 seed 数据含 source_note，AI 回答必须引用"),
    ("④ 空数据兜底", "无匹配时声明「数据层未收录该型号，请查阅官方手册」"),
]
for i, (title, desc) in enumerate(shield_data):
    y = Inches(5.3 + i * 0.42)
    rect(s, Inches(6.9), y, Inches(5.5), Inches(0.37), fill=DARK_CARD, corner=True)
    tb(s, Inches(7.05), y + Inches(0.03), Inches(5.2), Inches(0.32),
       f"{title}  —  {desc}", size=9.5, color=LIGHT_GRAY)

page_num(s, 4)

# ═══════════════════════════════════════════
# SLIDE 5 — TECHNICAL ARCHITECTURE + SCREENSHOTS
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "技术架构 & 产品实现", "Next.js 15 全栈 + Supabase PostgreSQL + DeepSeek API + Vercel AI SDK")

# 3-column architecture
arch = [
    ("🖥️ 前端层", TEAL, ["Next.js 15 App Router", "React Server + Client Components", "Tailwind CSS v4 + shadcn/ui", "useChat 流式消费 + react-markdown", "next-themes 暗色模式"]),
    ("⚙️ 服务层", STEEL, ["API Route Handlers (SSE/非流式)", "Auth Middleware (Cookie+Bearer)", "Supabase RLS 数据隔离", "Zod 输入校验", "crypto.randomUUID 预生成 ID"]),
    ("🗄️ 数据 & AI 层", RGBColor(0x00, 0x6B, 0x65), ["Supabase PostgreSQL (8 张表)", "EAV 属性模型 (product_attributes)", "SQL ILIKE 关键词检索", "DeepSeek Chat API (OpenAI兼容)", "DeepSeek Chat API (OpenAI兼容)"]),
]
for i, (title, accent, items) in enumerate(arch):
    x = Inches(0.7 + i * 4.1)
    y = Inches(1.7)
    card = rect(s, x, y, Inches(3.85), Inches(2.85), fill=DARK_CARD, corner=True)
    rect(s, x, y, Inches(3.85), Pt(4), fill=accent)
    tb(s, x + Inches(0.2), y + Inches(0.15), Inches(3.45), Inches(0.35),
       title, size=14, color=accent, bold=True)
    for j, item in enumerate(items):
        tb(s, x + Inches(0.3), y + Inches(0.7 + j * 0.36), Inches(3.35), Inches(0.35),
           f"▸ {item}", size=10, color=LIGHT_GRAY)

# Bottom: Product screenshots
tb(s, Inches(0.7), Inches(4.85), Inches(11.5), Inches(0.3),
   "📱  产品界面实拍", size=14, color=TEAL, bold=True)

add_screenshot(s, "02_chat.png", Inches(0.7), Inches(5.2), Inches(4.0), Inches(2.1))
add_screenshot(s, "03_products.png", Inches(4.9), Inches(5.2), Inches(4.0), Inches(2.1))
add_screenshot(s, "05_solutions_new.png", Inches(9.1), Inches(5.2), Inches(4.0), Inches(2.1))

# Labels under screenshots
for i, (label, x) in enumerate([("AI 对话选型", Inches(0.7)), ("产品数据浏览", Inches(4.9)), ("方案生成", Inches(9.1))]):
    tb(s, x + Inches(1.2), Inches(7.25), Inches(1.6), Inches(0.2),
       label, size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

page_num(s, 5)

# ═══════════════════════════════════════════
# SLIDE 6 — VERIFICATION & IMPACT
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "验证 & 成果", "端到端测试 → 准确率验证 → 效果量化")

# Metrics row
metrics = [
    ("5-15s", "AI 问答响应时间\n(传统流程 1-2 天)", GREEN),
    ("100%", "对话持久化成功率\n(onFinish 服务端写入)", TEAL),
    ("精确匹配", "型号参数来自真实 DB\n幻觉率 0%", STEEL),
    ("2 周", "0→1 MVP 交付\n独立全栈开发", AMBER),
]
for i, (num, label, accent) in enumerate(metrics):
    x = Inches(0.7 + i * 3.08)
    y = Inches(1.7)
    card = rect(s, x, y, Inches(2.85), Inches(1.3), fill=DARK_CARD, corner=True)
    rect(s, x, y, Inches(2.85), Pt(3), fill=accent)
    tb(s, x + Inches(0.15), y + Inches(0.08), Inches(2.55), Inches(0.55),
       num, size=28, color=accent, bold=True)
    tb(s, x + Inches(0.15), y + Inches(0.7), Inches(2.55), Inches(0.5),
       label, size=10, color=MID_GRAY, spacing=1.3)

# Left: Test case
rich(s, Inches(0.7), Inches(3.35), Inches(5.8), Inches(0.35), [
    {"text": "🧪  真实测试案例验证", "size": 15, "color": TEAL, "bold": True, "spacing": 6},
])

rect(s, Inches(0.7), Inches(3.75), Inches(5.8), Inches(3.3), fill=DARK_CARD, corner=True)
rich(s, Inches(0.9), Inches(3.85), Inches(5.4), Inches(3.1), [
    {"text": "👤 输入", "size": 12, "color": TEAL, "bold": True, "spacing": 4},
    {"text": "\"What is Schneider TM221CE40R PLC specs? I/O count and protocols.\"", "size": 10.5, "color": WHITE, "spacing": 6},
    {"text": "🤖 AI 输出验证 (全部通过 ✅)", "size": 12, "color": TEAL, "bold": True, "spacing": 4},
    {"text": "✅ 型号 TM221CE40R 精确匹配 seed 数据", "size": 10.5, "color": GREEN, "spacing": 2},
    {"text": "✅ I/O = 40 (24入16出) ← 来自 product_attributes", "size": 10.5, "color": GREEN, "spacing": 2},
    {"text": "✅ Modbus RTU/TCP ← 来自 product_attributes", "size": 10.5, "color": GREEN, "spacing": 2},
    {"text": "✅ 跨品牌替代方案：XINJE XD3-32R-E", "size": 10.5, "color": GREEN, "spacing": 2},
    {"text": "✅ 来源标注：「示例型号」「以官方资料为准」", "size": 10.5, "color": GREEN, "spacing": 2},
    {"text": "✅ 空查询时声明「数据层未收录，请查阅官方手册」", "size": 10.5, "color": GREEN, "spacing": 2},
])

# Right: Impact + screenshot
rich(s, Inches(6.8), Inches(3.35), Inches(5.8), Inches(0.35), [
    {"text": "📊  效果 & 产品价值", "size": 15, "color": TEAL, "bold": True, "spacing": 6},
])

rect(s, Inches(6.8), Inches(3.75), Inches(5.8), Inches(1.7), fill=DARK_CARD, corner=True)
rich(s, Inches(7.0), Inches(3.9), Inches(5.4), Inches(1.4), [
    {"text": "🔑 核心价值验证", "size": 13, "color": TEAL, "bold": True, "spacing": 6},
    {"text": "AI 可在 5 分钟内替代售前工程师 80% 的常见选型工作", "size": 12, "color": WHITE, "bold": True, "spacing": 4},
    {"text": "✧ 响应速度：1-2 天 → 5-15 秒（1000x 提升）", "size": 10.5, "color": LIGHT_GRAY, "spacing": 2},
    {"text": "✧ 选型准确率：基于数据库，杜绝编造", "size": 10.5, "color": LIGHT_GRAY, "spacing": 2},
    {"text": "✧ 知识沉淀：对话自动保存，新人可回溯学习", "size": 10.5, "color": LIGHT_GRAY, "spacing": 2},
    {"text": "✧ 跨品牌对比：施耐德 vs 信捷，AI 自动完成", "size": 10.5, "color": LIGHT_GRAY, "spacing": 2},
])

# Screenshot at bottom right
add_screenshot(s, "03_products.png", Inches(6.8), Inches(5.55), Inches(5.8), Inches(1.55))

page_num(s, 6)

# ═══════════════════════════════════════════
# SLIDE 7 — SUMMARY & ROADMAP
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "总结 & 路线图", "产品成果回顾  ·  AI PM 能力体现  ·  未来规划")

# Achievements
rect(s, Inches(0.7), Inches(1.7), Inches(7.5), Inches(3.8), fill=DARK_CARD, corner=True)
rich(s, Inches(0.9), Inches(1.85), Inches(7.1), Inches(3.5), [
    {"text": "🏆  项目核心成果", "size": 18, "color": TEAL, "bold": True, "spacing": 12},
    {"text": "▸ 0→1 完整 AI SaaS 产品：选型问答 + 产品浏览 + 方案生成三大闭环", "size": 12, "color": WHITE, "spacing": 4},
    {"text": "▸ RAG 检索增强架构：基于真实业务数据推荐，而非 LLM 通用知识", "size": 12, "color": WHITE, "spacing": 4},
    {"text": "▸ 生产级防幻觉设计：四层机制，端到端测试验证，幻觉率 0%", "size": 12, "color": WHITE, "spacing": 4},
    {"text": "▸ 流式 AI 体验：5-15 秒实时输出，用户即时看到回答", "size": 12, "color": WHITE, "spacing": 4},
    {"text": "▸ EAV 灵活数据模型 + RLS 安全隔离 + 服务端持久化，直接可用", "size": 12, "color": WHITE, "spacing": 4},
    {"text": "", "size": 6, "color": WHITE},
    {"text": "💪  体现的 AI PM 能力", "size": 18, "color": TEAL, "bold": True, "spacing": 12},
    {"text": "▸ 行业洞察 — 深入理解工业自动化分销链路，找到 AI 切入点", "size": 11, "color": LIGHT_GRAY, "spacing": 3},
    {"text": "▸ 产品决策 — MVP 范围取舍，知道「不做」比「做」更难", "size": 11, "color": LIGHT_GRAY, "spacing": 3},
    {"text": "▸ AI 工程化 — Prompt 工程 + RAG + 防幻觉 + 数据追溯，生产级设计", "size": 11, "color": LIGHT_GRAY, "spacing": 3},
    {"text": "▸ 技术落地 — Next.js + Supabase + AI SDK 全栈独立开发，2 周交付", "size": 11, "color": LIGHT_GRAY, "spacing": 3},
])

# Roadmap
rect(s, Inches(8.5), Inches(1.7), Inches(4.5), Inches(3.8), fill=DARK_CARD, corner=True)
rich(s, Inches(8.7), Inches(1.85), Inches(4.1), Inches(3.5), [
    {"text": "🗺️  产品路线图", "size": 18, "color": TEAL, "bold": True, "spacing": 14},
    {"text": "V1.1  短期", "size": 14, "color": AMBER, "bold": True, "spacing": 4},
    {"text": "组织协作（销售↔售前协同）\n审批流（大项目方案审批）\n方案生成优化", "size": 10.5, "color": LIGHT_GRAY, "spacing": 10},
    {"text": "V1.2 – V2.0  中期", "size": 14, "color": AMBER, "bold": True, "spacing": 4},
    {"text": "采购角色接入（货期/成本）\n微信小程序（移动端选型）\npgvector 向量知识库", "size": 10.5, "color": LIGHT_GRAY, "spacing": 10},
    {"text": "V3.0+  长期", "size": 14, "color": AMBER, "bold": True, "spacing": 4},
    {"text": "多品牌 (汇川/西门子/ABB)\n垂直行业模板\n数据分析看板", "size": 10.5, "color": LIGHT_GRAY, "spacing": 8},
])

# Bottom tagline
rect(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.9), fill=DARK_CARD, corner=True)
tb(s, Inches(1.0), Inches(5.98), Inches(11.3), Inches(0.7),
   '〝 让 AI 成为每一个工业销售的随身售前工程师 — 选型不再等待，知识不再流失 〞',
   size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

tb(s, Inches(0.7), Inches(7.0), Inches(11.9), Inches(0.3),
   "github.com/wangzibin/industrial-copilot  ·  Next.js 15 + Supabase + DeepSeek + Vercel AI SDK",
   size=9, color=DIM_GRAY, align=PP_ALIGN.CENTER)

page_num(s, 7)

# ── SAVE ──
prs.save(OUTPUT_PATH)
print(f"PPT saved to: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
